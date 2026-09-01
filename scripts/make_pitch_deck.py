#!/usr/bin/env python3
"""Generate a professional 16:9 pitch deck for client submission.

Outputs:
  out/Vehicle_Selector_Pro_Pitch_Deck.pptx  (editable)
  out/Vehicle_Selector_Pro_Pitch_Deck.pdf   (rendered, landscape)

Built entirely with free AI tooling.
"""
from __future__ import annotations

import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
# Override with REPORT_OUT=<dir> to generate a separate review copy without
# touching the published deliverables (which live in report/ and out/).
OUT = os.environ.get("REPORT_OUT", os.path.join(ROOT, "out"))
FRAMES = os.path.join(ROOT, "demo", "autoplay", "frames_live")
os.makedirs(OUT, exist_ok=True)

PPTX_PATH = os.path.join(OUT, "Vehicle_Selector_Pro_Pitch_Deck.pptx")
PDF_PATH = os.path.join(OUT, "Vehicle_Selector_Pro_Pitch_Deck.pdf")

ACCENT = "#E04B33"
DARK = "#282A2E"
GREY = "#6E6E6E"
LIGHT = "#FBEEE9"

REPO = "https://github.com/ai-dev-2024/vehicle-selector-pro"
DEMO = "https://vehicle-selector-pro.fly.dev/demo"
ADMIN = "https://vehicle-selector-pro.fly.dev/demo/admin"
APP = "https://vehicle-selector-pro.fly.dev"

# ---------------------------------------------------------------- deck content
SLIDES = [
    {
        "kind": "title",
        "kicker": "PITCH DECK · CLIENT PRESENTATION",
        "title": "Vehicle Selector Pro",
        "subtitle": "A Production-Grade Shopify App for Vehicle Fitment",
        "links": [("Live demo", DEMO), ("Merchant admin", ADMIN), ("Repository", REPO)],
    },
    {
        "kind": "screens2",
        "kicker": "THE PROBLEM",
        "title": "Aftermarket parts sell badly online — compatibility is the whole sale",
        "lead": "A brake pad or wiper blade only fits an exact Year / Make / Model / Trim / Engine. Wrong-fit means wrong parts, angry customers, and costly returns — so fitment data becomes the product.",
        "bullets": [
            "Dense compatibility tables bury the answer",
            "Per-product manual rules rot and cost hours",
            "Paid fitment feeds lock merchants into contracts",
            "Built end-to-end with free AI tooling on a 2018 ASUS ZenBook — proof of the headroom ahead",
        ],
        "shots": [("04_storefront_home.png", "Cascading vehicle selector — storefront"),
                  ("01_dashboard.png", "Polaris merchant command center")],
    },
    {
        "kind": "metrics",
        "kicker": "THE SOLUTION · ONE DATASET, SYNCED EVERYWHERE",
        "title": "A real, browsable demo catalog already in production",
        "metrics": [("48", "YMMTE configurations"), ("402", "verified fitments"),
                    ("74", "mapped products"), ("25", "brands represented")],
        "foot": "Fitments entered once into a normalized cache, pushed to Shopify metafields, and served to shoppers behind an HMAC-signed App Proxy — Guaranteed Exact / Does-NOT-Fit / Universal badges on every product page.",
    },
    {
        "kind": "arch",
        "kicker": "HOW IT WORKS",
        "title": "A cascading fitment engine, engineered for production",
    },
    {
        "kind": "cta",
        "kicker": "READY WHEN YOU ARE",
        "title": "See it live — watch it run",
        "lines": ["Try the live storefront demo", "Explore the merchant admin", "Open the repository", "Watch the two narrated walkthroughs"],
        "links": [DEMO, ADMIN, REPO, "https://ai-dev-2024.github.io/vehicle-selector-pro/demo/videos.html"],
    },
]


def shot(name):
    return os.path.join(FRAMES, name)


# ===========================================================================
# PDF deck (reportlab landscape)
# ===========================================================================
def build_pdf() -> None:
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER
    from reportlab.graphics.shapes import Drawing, Rect, String
    from reportlab.graphics import renderPDF
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
                                    PageBreak, Image as RLImage)
    from reportlab.platypus import HRFlowable

    W, H = 13.333 * inch, 7.5 * inch
    M = 0.6 * inch
    ACC = colors.HexColor(ACCENT); DARKC = colors.HexColor(DARK); GREYC = colors.HexColor(GREY)
    LIGHT_C = colors.HexColor(LIGHT)

    ss = getSampleStyleSheet()
    def PS(**kw):
        return ParagraphStyle(f"d{abs(hash(frozenset(kw.items())))%10**8}", parent=ss["BodyText"], **kw)

    def P(t, **kw): return Paragraph(t, PS(**kw))

    story = []

    def add_slide():
        story.append(PageBreak())

    def header(kicker, title, page_num):
        story.append(P(f"<font color='{ACCENT}'><b>{kicker}</b></font>",
                       fontSize=10, leading=12, spaceAfter=8))
        story.append(P(f"<b>{title}</b>", fontSize=26, leading=30, textColor=DARKC, spaceAfter=8))
        story.append(HRFlowable(width="100%", thickness=2, color=ACC, spaceAfter=16))

    def bullets(items):
        for it in items:
            story.append(P(f"<font color='{ACCENT}'>•</font>  {it}", fontSize=13, leading=20,
                           textColor=DARKC, leftIndent=6, spaceAfter=8))

    def two_cols(pairs, wfraction=0.48):
        rows = []
        for head, its in pairs:
            cell = []
            cell.append(P(f"<b><font color='{ACCENT}'>{head}</font></b>", fontSize=14, spaceAfter=6))
            for x in its:
                cell.append(P(f"•  {x}", fontSize=12, leading=18, textColor=DARKC, spaceAfter=4))
            cell.append(Spacer(1, 0.2 * inch))
            rows.append(cell)
        t = Table([rows], colWidths=[W - 2 * M, W - 2 * M])
        t.setStyle(TableStyle([
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 6), ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 4), ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        # split into per-column widths properly
        colW = (W - 2 * M - 24) / 2
        t = Table([rows[0]], colWidths=[colW, colW])
        t.setStyle(TableStyle([
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4), ("RIGHTPADDING", (0, 0), (-1, -1), 4),
        ]))
        return t

    # 1 Title
    story.append(P("P I T C H   D E C K", fontSize=11, leading=14, textColor=ACC,
                   fontName="Helvetica-Bold", spaceAfter=22))
    story.append(P("Vehicle Selector Pro", fontSize=44, leading=48, textColor=DARKC,
                   fontName="Helvetica-Bold", spaceAfter=8))
    story.append(P("A Production-Grade Shopify App for Vehicle Fitment", fontSize=19, leading=24,
                   textColor=DARKC, spaceAfter=6))
    story.append(HRFlowable(width="40%", thickness=3, color=ACC, spaceBefore=8, spaceAfter=22))
    story.append(P("CLIENT DELIVERABLE  ·  AUGUST 2026", fontSize=11, leading=14,
                   textColor=GREYC, spaceAfter=30))
    for l, u in SLIDES[0]["links"]:
        story.append(P(f"<b>{l}:</b> <a href='{u}'><font color='{ACCENT}'>{u}</font></a>",
                       fontSize=12, leading=18))

    page_num = 1
    for s in SLIDES[1:]:
        page_num += 1
        kind = s["kind"]
        from reportlab.platypus import PageBreak as _PB
        story.append(_PB())
        header(s["kicker"], s["title"], page_num)
        if kind == "content":
            if "cols" in s:
                story.append(two_cols(s["cols"]))
            else:
                bullets(s["bullets"])
        elif kind == "screens2":
            # Lead + bullets, then two proper-aspect screenshots sized to fit the remaining frame.
            story.append(P(s["lead"], fontSize=12, leading=17, textColor=DARKC, spaceAfter=6))
            for it in s["bullets"]:
                story.append(P(f"<font color='{ACCENT}'>•</font>  {it}", fontSize=11, leading=15,
                               textColor=DARKC, leftIndent=6, spaceAfter=3))
            # Cap image height to fit the remaining frame; derive width to keep native 16:9.
            sh = 2.5 * inch
            sw = sh * (1920 / 1080)
            cells = []
            for name, cap_txt in s["shots"]:
                cell = []
                p = shot(name)
                if os.path.exists(p):
                    cell.append(RLImage(p, width=sw, height=sh))
                cell.append(Spacer(1, 1))
                cell.append(P(cap_txt, fontSize=8, leading=10, textColor=GREYC, alignment=1))
                cells.append(cell)
            t = Table([cells], colWidths=[(W - 2 * M) / 2, (W - 2 * M) / 2])
            t.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"),
                                   ("LEFTPADDING", (0, 0), (-1, -1), 6),
                                   ("RIGHTPADDING", (0, 0), (-1, -1), 6)]))
            story.append(t)
        elif kind == "metrics":
            data = [[""] * 4]
            heads = [Paragraph(f"<font color='{ACCENT}'><font size='28'><b>{m[0]}</b></font></font><br/><font size='12'>{m[1]}</font>",
                     PS(alignment=TA_CENTER, leading=18)) for m in s["metrics"]]
            t = Table([heads], colWidths=[(W - 2 * M) / 4] * 4)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), LIGHT_C),
                ("GRID", (0, 0), (-1, -1), 0.6, ACC),
                ("TOPPADDING", (0, 0), (-1, -1), 30), ("BOTTOMPADDING", (0, 0), (-1, -1), 30),
                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
            ]))
            story.append(Spacer(1, 0.8 * inch))
            story.append(t)
            story.append(Spacer(1, 0.5 * inch))
            story.append(P(s.get("foot", "Try it live — no setup required. Every number above is visible in the running app."),
                           fontSize=12, textColor=GREYC))
        elif kind == "arch":
            from reportlab.graphics.shapes import Line, Polygon
            from reportlab.platypus import Flowable as _Flowable

            class Arch(_Flowable):
                def __init__(self, w, h):
                    super().__init__(); self.w = w; self.h = h

                def wrap(self, aw, ah):
                    return self.w, self.h

                def draw(self):
                    d = Drawing(self.w, self.h)
                    nodes = [("sf", "Shopify Storefront", "Theme App Extension"),
                             ("ap", "App Proxy", "HMAC JSON"),
                             ("ct", "Controllers", "admin · webhooks"),
                             ("jb", "Sidekiq / Redis", "async jobs"),
                             ("db", "PostgreSQL", "cache"),
                             ("gql", "Shopify GraphQL", "metafieldsSet")]
                    nw = (self.w - 0.4 * inch) / (len(nodes) - 1)
                    bw = min(1.5 * inch, nw - 0.18 * inch)
                    bh = 1.6 * inch
                    y0 = (self.h - bh) / 2
                    xs = [0.2 * inch + i * nw for i in range(len(nodes))]
                    centers = {}
                    for i, (nid, label, sub) in enumerate(nodes):
                        x = xs[i] - bw / 2; centers[nid] = (xs[i], y0 + bh / 2)
                        d.add(Rect(x, y0, bw, bh, rx=8, ry=8, fillColor=LIGHT_C, strokeColor=ACC,
                                   strokeWidth=1.2))
                        d.add(String(xs[i], y0 + bh - 22, label, textAnchor="middle",
                                     fontName="Helvetica-Bold", fontSize=10, fillColor=DARKC))
                        d.add(String(xs[i], y0 + 12, sub, textAnchor="middle", fontName="Helvetica",
                                     fontSize=7.5, fillColor=GREYC))
                    for a, b in [("sf", "ap"), ("sf", "ct"), ("ap", "db"), ("ct", "db"),
                                 ("ct", "jb"), ("jb", "gql"), ("gql", "sf"), ("jb", "db")]:
                        x1, _ = centers[a]; x2, cy = centers[b]
                        d.add(Line(x1 + bw / 2, cy, x2 - bw / 2, cy, strokeColor=ACC, strokeWidth=1.4))
                        dx = 1 if x2 >= x1 else -1
                        tip = x2 - (bw / 2) * dx
                        d.add(Polygon([tip, cy, tip - 5 * dx, cy - 3, tip - 5 * dx, cy + 3],
                                      fillColor=ACC, strokeColor=ACC))
                    renderPDF.draw(d, self.canv, 0, 0)

            from reportlab.platypus import Spacer as Sp
            story.append(Sp(1, 0.5 * inch))
            story.append(Arch(W - 2 * M, 3.3 * inch))
            story.append(Sp(1, 0.4 * inch))
            story.append(P("Zero ScriptTag API — pure Theme App Extension, signed App Proxy, and safe webhooks.",
                           fontSize=12, textColor=GREYC))
        elif kind == "screens":
            colW = (W - 2 * M - 30) / 2
            cell_pairs = []
            imgs = s["shots"]
            left = imgs[:2]; right = imgs[2:]
            def cell(imgs):
                out = []
                for name, cap_txt in imgs:
                    p = shot(name)
                    if not os.path.exists(p): continue
                    iw = 3.4 * inch
                    ih = iw * (1080 / 1920)
                    img = RLImage(p, width=iw, height=ih)
                    img.hAlign = "CENTER"
                    out.append(img)
                    out.append(Spacer(1, 2))
                    out.append(P(cap_txt, fontSize=9, leading=11, textColor=GREYC,
                                 alignment=1, spaceAfter=8))
                return out
            t = Table([[cell(left), cell(right)]], colWidths=[colW, colW])
            t.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"),
                                   ("LEFTPADDING", (0, 0), (-1, -1), 4),
                                   ("RIGHTPADDING", (0, 0), (-1, -1), 4)]))
            story.append(t)
        elif kind == "statement":
            story.append(Spacer(1, 0.7 * inch))
            t = Table([[P(s["text"], fontSize=16, leading=24, textColor=DARKC)]],
                      colWidths=[9.0 * inch])
            t.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), LIGHT_C),
                                   ("BOX", (0, 0), (-1, -1), 1.5, ACC),
                                   ("TOPPADDING", (0, 0), (-1, -1), 24),
                                   ("BOTTOMPADDING", (0, 0), (-1, -1), 24),
                                   ("LEFTPADDING", (0, 0), (-1, -1), 24),
                                   ("RIGHTPADDING", (0, 0), (-1, -1), 24)]))
            t.hAlign = "LEFT"
            story.append(t)
        elif kind == "cta":
            story.append(Spacer(1, 0.4 * inch))
            for l, u in zip(s["lines"], s["links"]):
                story.append(P(f"<font color='{ACCENT}'>›</font>  <a href='{u}'><font color='{DARKC}'><b>{l}</b></font></a>  —  <font color='{GREYC}'>{u}</font>",
                               fontSize=15, leading=30))
            story.append(Spacer(1, 0.5 * inch))
            story.append(P("Both the live storefront and merchant admin are running right now.",
                           fontSize=12, textColor=GREYC))

        # slide footer
        story.append(Spacer(1, 0.25 * inch))
        story.append(P(f"{page_num:02d}", fontSize=9, textColor=GREYC))

    def on_page(canvas, doc):
        canvas.saveState()
        canvas.setFillColor(ACC)
        canvas.rect(0, H - 0.14 * inch, W, 0.09 * inch, fill=1, stroke=0)
        canvas.setFillColor(GREYC); canvas.setFont("Helvetica", 8)
        canvas.drawString(M, 0.3 * inch, "Vehicle Selector Pro — Pitch Deck")
        canvas.drawRightString(W - M, 0.3 * inch, "vehicle-selector-pro.fly.dev")
        canvas.restoreState()

    doc = SimpleDocTemplate(PDF_PATH, pagesize=(W, H), leftMargin=M, rightMargin=M,
                            topMargin=0.7 * inch, bottomMargin=0.5 * inch,
                            title="Vehicle Selector Pro — Pitch Deck", author="Vehicle Selector Pro")
    doc.build(story, onFirstPage=on_page, onLaterPages=on_page)
    if os.path.exists(os.path.join(OUT, "_arch_tmp.png")):
        os.remove(os.path.join(OUT, "_arch_tmp.png"))
    print("WROTE", PDF_PATH)


# ===========================================================================
# PPTX deck
# ===========================================================================
def build_pptx() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    ACC = RGBColor(0xE0, 0x4B, 0x33)
    DARKC = RGBColor(0x28, 0x2A, 0x2E)
    GREYC = RGBColor(0x6E, 0x6E, 0x6E)
    LIGHT_C = RGBColor(0xFB, 0xEE, 0xE9)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    W, H = 13.333, 7.5

    def add_rect(slide, l, t, w, h, fill, line=ACC, text="", bold=False, font=14, color=DARKC,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_w=None, shape=MSO_SHAPE.RECTANGLE):
        sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            if line_w: sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(font); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"
        return sp

    def title_slide():
        s = prs.slides.add_slide(BLANK)
        add_rect(s, 0, 0, W, 0.14, ACC)
        add_rect(s, 0, H - 0.14, W, 0.14, ACC)
        add_rect(s, 0.9, 1.1, 0.06, 3.6, ACC)
        add_rect(s, 1.2, 1.35, 9, 0.5, None, None, "PITCH DECK  ·  CLIENT PRESENTATION",
                 bold=True, font=12, color=ACC)
        add_rect(s, 1.2, 2.0, 11, 1.2, None, None, "Vehicle Selector Pro", bold=True, font=44, color=DARKC)
        add_rect(s, 1.2, 3.15, 10.5, 0.8, None, None,
                 "A Production-Grade Shopify App for Vehicle Fitment", font=19, color=DARKC)
        add_rect(s, 1.25, 4.2, 3.2, 0.06, ACC)
        add_rect(s, 1.2, 4.5, 10, 0.4, None, None, "CLIENT DELIVERABLE  ·  AUGUST 2026",
                 font=11, color=GREYC)
        y = 5.1
        for label, url in SLIDES[0]["links"]:
            tb = s.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11), Inches(0.4))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            rl = p.add_run(); rl.text = f"{label}:  "; rl.font.size = Pt(13)
            rl.font.bold = True; rl.font.color.rgb = DARKC; rl.font.name = "Calibri"
            ru = p.add_run(); ru.text = url; ru.font.size = Pt(13)
            ru.font.color.rgb = ACC; ru.font.name = "Calibri"; ru.font.underline = True
            ru.hyperlink.address = url
            y += 0.5

    def content_slide(kick, title):
        s = prs.slides.add_slide(BLANK)
        add_rect(s, 0, 0, W, 0.14, ACC)
        add_rect(s, 0.9, 0.6, 6, 0.4, None, None, kick.upper(), bold=True, font=12, color=ACC)
        add_rect(s, 0.9, 1.0, 11.5, 1.0, None, None, title, bold=True, font=28, color=DARKC)
        add_rect(s, 0.9, 2.0, 11.5, 0.03, ACC)
        return s

    def bullets_box(slide, items, l=1.0, t=2.3, w=11.2, h=4.4, size=15, leading=1.9, color=DARKC):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for (head, body) in [(None, i) for i in items]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run(); r.text = "▪  "
            r.font.color.rgb = ACC; r.font.size = Pt(size); r.font.bold = True
            r2 = p.add_run(); r2.text = body
            r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Calibri"
            p.space_after = Pt(leading * 4)
            p.line_spacing = leading
        return tb

    def two_col_box(slide, pairs, l=1.0, t=2.4, w=11.2, h=4.2):
        cw = (w - 0.5) / 2
        for i, (head, its) in enumerate(pairs):
            x = l + i * (cw + 0.5)
            add_rect(slide, x, t, cw, h, LIGHT_C, ACC, head, bold=True, font=18, color=ACC,
                     anchor=MSO_ANCHOR.TOP)
            tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(t + 0.7), Inches(cw - 0.5),
                                          Inches(h - 0.9))
            tf = tb.text_frame; tf.word_wrap = True; first = True
            for it in its:
                p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
                rr = p.add_run(); rr.text = f"▪  {it}"
                rr.font.size = Pt(13); rr.font.color.rgb = DARKC; rr.font.name = "Calibri"
                p.space_after = Pt(10)

    def metrics_slide(slide, s):
        cw = (11.2 - 0.9) / 4
        for i, (num, label) in enumerate(s["metrics"]):
            x = 1.0 + i * (cw + 0.3)
            add_rect(slide, x, 2.9, cw, 2.4, LIGHT_C, ACC, num, bold=True, font=34, color=ACC,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_rect(slide, x, 4.6, cw, 0.9, None, None, label, font=13, color=DARKC,
                     align=PP_ALIGN.CENTER)

    title_slide()
    for s in SLIDES[1:]:
        kind = s["kind"]
        if kind == "content":
            sl = content_slide(s["kicker"], s["title"])
            if "cols" in s:
                two_col_box(sl, s["cols"])
            else:
                bullets_box(sl, s["bullets"])
        elif kind == "screens2":
            sl = content_slide(s["kicker"], s["title"])
            tb = sl.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.2), Inches(1.1))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = s["lead"]
            r.font.size = Pt(13); r.font.color.rgb = GREYC; r.font.name = "Calibri"
            bb = bullets_box(sl, s["bullets"], l=1.0, t=2.9, w=11.2, h=1.3, size=12, leading=1.5)
            cw = (11.2 - 0.6) / 2; colw = (cw - 0.25) / 2
            for idx, (name, capt) in enumerate(s["shots"]):
                col = idx % 2; row = int(idx / 2)
                x = 1.0 + col * (colw + 0.25)
                y = 4.0 + row * (1.6)
                p = shot(name)
                if not os.path.exists(p): continue
                from PIL import Image
                sl.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(colw - 0.1), height=Inches((colw - 0.1) * 0.5625))
        elif kind == "metrics":
            sl = content_slide(s["kicker"], s["title"]); metrics_slide(sl, s)
            tb = sl.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.2), Inches(1.0))
            tf = tb.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = s.get("foot", "")
            r.font.size = Pt(12); r.font.color.rgb = GREYC; r.font.name = "Calibri"
        elif kind == "arch":
            sl = content_slide(s["kicker"], s["title"])
            nodes = ["Storefront", "App Proxy", "Controllers", "Sidekiq/Redis", "PostgreSQL", "GraphQL"]
            nw = (11.2 - 1.) / (len(nodes) - 1)
            bw = 1.5
            y = 3.4
            for i, lab in enumerate(nodes):
                x = 1.0 + i * nw
                add_rect(sl, x - bw / 2, y, bw, 1.5, LIGHT_C, ACC, lab, bold=True, font=13,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_rect(sl, 1.0, 5.4, 11, 0.5, None, None,
                     "Zero ScriptTag API — pure Theme App Extension, signed App Proxy, safe webhooks.",
                     font=13, color=GREYC)
        elif kind == "screens":
            sl = content_slide(s["kicker"], s["title"])
            cw = (11.2 - 0.6) / 2; colw = (cw - 0.25) / 2
            imgs = s["shots"]
            for idx, (name, capt) in enumerate(imgs):
                col = idx % 2; row = int(idx / 2)
                x = 1.0 + col * (colw + 0.2 + cw) - cw if col else 1.0
                x = 1.0 + col * (colw + 0.25)
                y = 2.3 + row * (2.3)
                p = shot(name)
                if not os.path.exists(p): continue
                from PIL import Image
                iw, ih = Image.open(p).size
                pw = colw - 0.1
                ph = pw * ih / iw
                if row == 0: pass
                sl.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(colw - 0.1), height=Inches(ph))
                add_rect(sl, x, y + ph + 0.02, colw - 0.1, 0.35, None, None, capt, font=10,
                         color=GREYC, align=PP_ALIGN.CENTER)
        elif kind == "statement":
            sl = content_slide(s["kicker"], s["title"])
            add_rect(sl, 1.0, 2.6, 11, 3.0, LIGHT_C, ACC, s["text"], font=18, color=DARKC,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif kind == "cta":
            sl = content_slide(s["kicker"], s["title"])
            y = 2.8
            for lab, url in zip(s["lines"], s["links"]):
                add_rect(sl, 1.0, y, 0.4, 0.4, ACC, None, "›", bold=True, font=18, color=RGBColor(255,255,255),
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                lab_box = add_rect(sl, 1.5, y - 0.05, 11, 0.5, None, None, lab,
                                   bold=True, font=16, color=DARKC)
                try:
                    lab_box.text_frame.paragraphs[0].runs[0].hyperlink.address = url
                except Exception:
                    pass
                y += 0.8
            add_rect(sl, 1.0, y + 0.2, 11, 0.4, None, None,
                     "Both the live storefront and merchant admin are running right now.",
                     font=12, color=GREYC)

    prs.save(PPTX_PATH)
    print("WROTE", PPTX_PATH)


if __name__ == "__main__":
    build_pdf()
    build_pptx()
    print("Done. Deck in:", OUT)